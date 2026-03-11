from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional
from urllib.error import HTTPError, URLError
from urllib.parse import urlencode
from urllib.request import Request, urlopen

import requests
from bs4 import BeautifulSoup

from report_generator.models import ExtractedDocument

DEFAULT_USER_AGENT = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0 Safari/537.36"


def _build_result(
    path: Path,
    *,
    parser_name: str,
    status: str,
    text: str = "",
    metadata: Optional[Dict[str, Any]] = None,
    warnings: Optional[List[str]] = None,
    error_code: Optional[str] = None,
    error_message: Optional[str] = None,
) -> ExtractedDocument:
    return ExtractedDocument(
        source_path=str(path),
        file_name=path.name,
        extension=path.suffix.lower(),
        parser_name=parser_name,
        status=status,
        text=text,
        metadata=metadata or {},
        warnings=warnings or [],
        error_code=error_code,
        error_message=error_message,
    )


def parse_text_file(path: Path) -> ExtractedDocument:
    encodings = ["utf-8", "utf-8-sig", "gbk", "gb18030"]
    for encoding in encodings:
        try:
            text = path.read_text(encoding=encoding)
            warnings = []
            if not text.strip():
                warnings.append("文件内容为空")
            return _build_result(
                path,
                parser_name="text_reader",
                status="ok",
                text=text,
                metadata={"encoding": encoding},
                warnings=warnings,
            )
        except UnicodeDecodeError:
            continue
        except OSError as exc:
            return _build_result(
                path,
                parser_name="text_reader",
                status="error",
                error_code="read_error",
                error_message=str(exc),
            )

    return _build_result(
        path,
        parser_name="text_reader",
        status="error",
        error_code="decode_error",
        error_message="无法使用预设编码读取文本文件",
    )


def parse_pdf_file(path: Path) -> ExtractedDocument:
    from pypdf import PdfReader

    try:
        reader = PdfReader(str(path))
        pages: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            content = page.extract_text() or ""
            pages.append(f"[第{index}页]\n{content}")
        text = "\n\n".join(pages)
        warnings = []
        if not text.strip():
            warnings.append("PDF 未提取到文本，可能是扫描件或加密文档")
        return _build_result(
            path,
            parser_name="pdf_reader",
            status="ok",
            text=text,
            metadata={"page_count": len(reader.pages)},
            warnings=warnings,
        )
    except Exception as exc:  # noqa: BLE001
        return _build_result(
            path,
            parser_name="pdf_reader",
            status="error",
            error_code="pdf_parse_error",
            error_message=str(exc),
        )


def parse_docx_file(path: Path) -> ExtractedDocument:
    from docx import Document

    try:
        document = Document(str(path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        tables: list[str] = []
        for table in document.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                tables.append("\n".join(rows))
        parts = []
        if paragraphs:
            parts.append("\n".join(paragraphs))
        if tables:
            parts.append("\n\n".join(tables))
        text = "\n\n".join(parts)
        warnings = []
        if not text.strip():
            warnings.append("Word 文档未提取到文本")
        return _build_result(
            path,
            parser_name="docx_reader",
            status="ok",
            text=text,
            metadata={"paragraph_count": len(paragraphs), "table_count": len(document.tables)},
            warnings=warnings,
        )
    except Exception as exc:  # noqa: BLE001
        return _build_result(
            path,
            parser_name="docx_reader",
            status="error",
            error_code="docx_parse_error",
            error_message=str(exc),
        )


def parse_pptx_file(path: Path) -> ExtractedDocument:
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
        slide_texts: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    texts.append(str(shape.text).strip())
            slide_texts.append(f"[第{index}页幻灯片]\n" + "\n".join(texts))
        text = "\n\n".join(slide_texts)
        warnings = []
        if not text.strip():
            warnings.append("PPT 未提取到文本")
        return _build_result(
            path,
            parser_name="pptx_reader",
            status="ok",
            text=text,
            metadata={"slide_count": len(presentation.slides)},
            warnings=warnings,
        )
    except Exception as exc:  # noqa: BLE001
        return _build_result(
            path,
            parser_name="pptx_reader",
            status="error",
            error_code="pptx_parse_error",
            error_message=str(exc),
        )


def parse_xlsx_file(path: Path) -> ExtractedDocument:
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(filename=str(path), data_only=True)
        sheet_blocks: list[str] = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if cell is None else str(cell).strip() for cell in row]
                if any(values):
                    rows.append(" | ".join(values))
            block = f"[工作表:{sheet.title}]\n" + "\n".join(rows)
            sheet_blocks.append(block)
        text = "\n\n".join(sheet_blocks)
        warnings = []
        if not text.strip():
            warnings.append("Excel 未提取到有效文本")
        return _build_result(
            path,
            parser_name="xlsx_reader",
            status="ok",
            text=text,
            metadata={"sheet_names": workbook.sheetnames},
            warnings=warnings,
        )
    except Exception as exc:  # noqa: BLE001
        return _build_result(
            path,
            parser_name="xlsx_reader",
            status="error",
            error_code="xlsx_parse_error",
            error_message=str(exc),
        )


def parse_file(path: Path) -> ExtractedDocument:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return parse_text_file(path)
    if suffix == ".pdf":
        return parse_pdf_file(path)
    if suffix == ".docx":
        return parse_docx_file(path)
    if suffix == ".pptx":
        return parse_pptx_file(path)
    if suffix == ".xlsx":
        return parse_xlsx_file(path)

    return _build_result(
        path,
        parser_name="unsupported",
        status="error",
        error_code="unsupported_extension",
        error_message=f"不支持的文件类型: {suffix}",
    )


def _extract_web_page_text(html: str) -> str:
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return re.sub(r"\n{2,}", "\n\n", text).strip()


def fetch_web_page(url: str, *, timeout: int = 15) -> dict[str, Any]:
    if os.getenv("REPORT_GENERATOR_ENABLE_WEB") != "1":
        return {
            "status": "disabled",
            "url": url,
            "title": "",
            "text": "",
            "warnings": [
                "当前环境未启用网页抓取。",
                "如需启用，请显式设置 REPORT_GENERATOR_ENABLE_WEB=1，并仅将网页内容作为补充线索使用。",
            ],
        }

    try:
        response = requests.get(url, timeout=timeout, headers={"User-Agent": DEFAULT_USER_AGENT})
        response.raise_for_status()
    except requests.RequestException as exc:
        return {
            "status": "error",
            "url": url,
            "title": "",
            "text": "",
            "warnings": [f"网页抓取失败: {exc}"],
        }

    html = response.text
    soup = BeautifulSoup(html, "html.parser")
    title = (soup.title.string or "").strip() if soup.title and soup.title.string else ""
    return {
        "status": "ok",
        "url": url,
        "title": title,
        "text": _extract_web_page_text(html),
        "warnings": [
            "网页正文抽取可能包含导航、页脚或转码噪声。",
            "写入报告前必须核验来源站点、发布日期、发布主体和原始上下文。",
        ],
    }


def search_web(query: str, *, limit: int = 5) -> dict[str, Any]:
    if os.getenv("REPORT_GENERATOR_ENABLE_WEB") != "1":
        return {
            "status": "disabled",
            "query": query,
            "results": [],
            "warnings": [
                "当前环境未启用联网搜索。",
                "如需启用，请显式设置 REPORT_GENERATOR_ENABLE_WEB=1；联网结果只能作为补充线索使用。",
            ],
        }

    search_url = os.getenv("REPORT_GENERATOR_WEB_SEARCH_URL", "https://www.baidu.com/s")
    selector = os.getenv("REPORT_GENERATOR_WEB_RESULT_SELECTOR", "h3 a")

    try:
        response = requests.get(
            search_url,
            params={"wd": query},
            timeout=15,
            headers={"User-Agent": DEFAULT_USER_AGENT},
        )
        response.raise_for_status()
    except requests.RequestException as exc:
        return {
            "status": "error",
            "query": query,
            "results": [],
            "warnings": [f"联网搜索请求失败: {exc}"],
        }

    soup = BeautifulSoup(response.text, "html.parser")
    candidates = soup.select(selector)
    results = []
    seen = set()

    for node in candidates:
        href = node.get("href", "").strip()
        title = node.get_text(" ", strip=True)
        if not href or not title or href in seen:
            continue
        seen.add(href)
        results.append(
            {
                "title": title,
                "url": href,
                "snippet": "",
                "source": search_url,
                "published_at": "",
                "retrieved_at": "",
            }
        )
        if len(results) >= limit:
            break

    return {
        "status": "ok",
        "query": query,
        "results": results,
        "warnings": [
            "当前为免费网页搜索抓取方案，不保证搜索页结构长期稳定。",
            "所有关键事实、数字、政策名称和发布日期，必须回到权威原始页面二次核验。",
            "建议优先采用政府官网、标准发布页、监管机构页面和官方公告。",
        ],
    }
